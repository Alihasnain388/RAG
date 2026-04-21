import hashlib
import os
import re
import tempfile
import uuid
from io import BytesIO
from pathlib import Path

import streamlit as st
from docx import Document as DocxDocument
from dotenv import load_dotenv
from langchain_core.documents import Document
from langchain_groq import ChatGroq
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_pinecone import PineconeVectorStore
from langchain_community.document_loaders import PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from openpyxl import load_workbook
from pinecone import Pinecone
from pptx import Presentation
from rank_bm25 import BM25Okapi
from sentence_transformers import CrossEncoder


load_dotenv()

st.set_page_config(page_title="RAG Chatbot")
st.title("RAG Chatbot")


def get_config_value(name: str) -> str | None:
    value = os.getenv(name)
    if value:
        return value

    try:
        value = st.secrets.get(name)
    except Exception:
        return None

    return value or None


PINECONE_API_KEY = get_config_value("PINECONE_API_KEY")
INDEX_NAME = get_config_value("PINECONE_INDEX_NAME")
GROQ_API_KEY = get_config_value("GROQ_API_KEY")

missing_config = [
    name
    for name, value in {
        "PINECONE_API_KEY": PINECONE_API_KEY,
        "PINECONE_INDEX_NAME": INDEX_NAME,
        "GROQ_API_KEY": GROQ_API_KEY,
    }.items()
    if not value
]

if missing_config:
    st.error("Missing required configuration: " + ", ".join(missing_config))
    st.stop()


@st.cache_resource(show_spinner=False)
def get_pinecone_index(api_key: str, index_name: str):
    pc = Pinecone(api_key=api_key)
    return pc.Index(index_name)


@st.cache_resource(show_spinner=False)
def get_embedding_model():
    return HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")


@st.cache_resource(show_spinner=False)
def get_reranker():
    return CrossEncoder("cross-encoder/ms-marco-MiniLM-L-6-v2")


@st.cache_resource(show_spinner=False)
def get_llm(api_key: str):
    return ChatGroq(model="llama-3.1-8b-instant", api_key=api_key)


def ensure_session_id() -> str:
    if "session_id" not in st.session_state:
        st.session_state.session_id = uuid.uuid4().hex
    return st.session_state.session_id


def uploaded_files_signature(uploaded_files) -> str:
    hasher = hashlib.sha256()

    for uploaded_file in uploaded_files:
        file_bytes = uploaded_file.getvalue()
        hasher.update(uploaded_file.name.encode("utf-8", errors="ignore"))
        hasher.update(str(len(file_bytes)).encode("utf-8"))
        hasher.update(hashlib.sha256(file_bytes).digest())

    return hasher.hexdigest()


def extract_pdf(uploaded_file) -> list[Document]:
    source_name = uploaded_file.name
    tmp_path = None

    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
            tmp.write(uploaded_file.getvalue())
            tmp_path = tmp.name

        loader = PyPDFLoader(tmp_path)
        docs = loader.load()

        for doc in docs:
            doc.metadata["source"] = source_name
            if "page" in doc.metadata:
                doc.metadata["page"] = doc.metadata["page"] + 1

        return docs
    finally:
        if tmp_path:
            try:
                os.unlink(tmp_path)
            except OSError:
                pass


def extract_docx(uploaded_file) -> list[Document]:
    docx_doc = DocxDocument(BytesIO(uploaded_file.getvalue()))
    parts = [paragraph.text for paragraph in docx_doc.paragraphs if paragraph.text.strip()]

    for table in docx_doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                parts.append(row_text)

    return [
        Document(
            page_content="\n".join(parts),
            metadata={"source": uploaded_file.name},
        )
    ]


def extract_pptx(uploaded_file) -> list[Document]:
    presentation = Presentation(BytesIO(uploaded_file.getvalue()))
    docs = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        parts = []

        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                parts.append(shape.text.strip())

            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        parts.append(row_text)

        if getattr(slide, "has_notes_slide", False):
            notes_frame = getattr(slide.notes_slide, "notes_text_frame", None)
            notes_text = notes_frame.text.strip() if notes_frame and notes_frame.text else ""
            if notes_text:
                parts.append("Speaker notes:\n" + notes_text)

        slide_text = "\n".join(parts).strip()
        if slide_text:
            docs.append(
                Document(
                    page_content=slide_text,
                    metadata={"source": uploaded_file.name, "slide": slide_number},
                )
            )

    return docs


def clean_cell(value) -> str:
    if value is None:
        return ""
    return str(value).strip()


def normalize_header(header: list[str]) -> list[str]:
    normalized = []

    for index, column in enumerate(header, start=1):
        normalized.append(column or f"Column {index}")

    return normalized


def find_column(headers: list[str], *candidates: str) -> str | None:
    header_lookup = {header.strip().lower(): header for header in headers}

    for candidate in candidates:
        match = header_lookup.get(candidate.strip().lower())
        if match:
            return match

    return None


def count_values(records: list[dict[str, str]], column: str) -> dict[str, int]:
    counts = {}

    for record in records:
        value = record.get(column, "").strip()
        if value:
            counts[value] = counts.get(value, 0) + 1

    return counts


def ordered_unique(values: list[str]) -> list[str]:
    unique = []
    seen = set()

    for value in values:
        value = value.strip()
        key = value.lower()
        if value and key not in seen:
            unique.append(value)
            seen.add(key)

    return unique


def excel_summary_docs(
    source_name: str,
    sheet_name: str,
    headers: list[str],
    records: list[dict[str, str]],
) -> list[Document]:
    if not records:
        return []

    docs = []
    client_column = find_column(headers, "Client")
    vm_column = find_column(headers, "VM Name", "VM", "Machine", "Machine Name")
    service_column = find_column(headers, "Service Name", "Service", "Services")
    project_column = find_column(headers, "Project")
    environment_column = find_column(headers, "Environment", "Env")

    lines = [
        f"Excel Summary for {source_name}",
        f"Sheet: {sheet_name}",
        f"Total data rows: {len(records)}",
        f"Columns: {', '.join(headers)}",
    ]

    if vm_column:
        unique_vms = ordered_unique([record.get(vm_column, "") for record in records])
        lines.append(f"Total unique VMs: {len(unique_vms)}")
        lines.append("Unique VM Names: " + ", ".join(unique_vms))

    if service_column:
        unique_services = ordered_unique(
            [record.get(service_column, "") for record in records]
        )
        lines.append(f"Total unique services: {len(unique_services)}")
        lines.append("Unique Services: " + ", ".join(unique_services))

    for label, column in [
        ("Client", client_column),
        ("Project", project_column),
        ("Environment", environment_column),
    ]:
        if not column:
            continue

        counts = count_values(records, column)
        if counts:
            formatted_counts = ", ".join(
                f"{value}: {count}" for value, count in sorted(counts.items())
            )
            lines.append(f"{label} counts: {formatted_counts}")

    docs.append(
        Document(
            page_content="\n".join(lines),
            metadata={
                "source": source_name,
                "sheet": sheet_name,
                "section": "Excel summary",
                "preserve_chunk": True,
            },
        )
    )

    if client_column and vm_column:
        clients = ordered_unique([record.get(client_column, "") for record in records])

        for client in clients:
            client_records = [
                record for record in records if record.get(client_column, "").strip() == client
            ]
            unique_vms = ordered_unique(
                [record.get(vm_column, "") for record in client_records]
            )
            unique_services = (
                ordered_unique([record.get(service_column, "") for record in client_records])
                if service_column
                else []
            )
            project_counts = (
                count_values(client_records, project_column) if project_column else {}
            )
            environment_counts = (
                count_values(client_records, environment_column)
                if environment_column
                else {}
            )

            client_lines = [
                f"Excel Client Summary for {client}",
                f"Source: {source_name}",
                f"Sheet: {sheet_name}",
                f"Client: {client}",
                f"Total rows for client: {len(client_records)}",
                f"Total unique VMs for client: {len(unique_vms)}",
                "Unique VM Names: " + ", ".join(unique_vms),
            ]

            if service_column:
                client_lines.extend(
                    [
                        f"Total unique services for client: {len(unique_services)}",
                        "Unique Services: " + ", ".join(unique_services),
                    ]
                )

            if project_counts:
                client_lines.append(
                    "Project counts: "
                    + ", ".join(
                        f"{value}: {count}" for value, count in sorted(project_counts.items())
                    )
                )
            if environment_counts:
                client_lines.append(
                    "Environment counts: "
                    + ", ".join(
                        f"{value}: {count}"
                        for value, count in sorted(environment_counts.items())
                    )
                )

            docs.append(
                Document(
                    page_content="\n".join(client_lines),
                    metadata={
                        "source": source_name,
                        "sheet": sheet_name,
                        "section": f"Client summary: {client}",
                        "client": client,
                        "preserve_chunk": True,
                    },
                )
            )

    return docs


def extract_excel(uploaded_file) -> list[Document]:
    workbook = load_workbook(BytesIO(uploaded_file.getvalue()), data_only=True, read_only=True)
    docs = []

    for sheet in workbook.worksheets:
        header = None
        records = []
        row_docs = []

        for row_number, row in enumerate(sheet.iter_rows(values_only=True), start=1):
            values = [clean_cell(value) for value in row]
            if not any(values):
                continue

            if header is None:
                header = normalize_header(values)
                continue

            padded_values = values + [""] * max(0, len(header) - len(values))
            row_map = {
                column: value
                for column, value in zip(header, padded_values)
                if value
            }
            records.append(row_map)

            row_text = "\n".join(
                f"- {column}: {value}" for column, value in row_map.items()
            )

            row_docs.append(
                Document(
                    page_content=f"Sheet: {sheet.title}\nRow: {row_number}\n{row_text}",
                    metadata={
                        "source": uploaded_file.name,
                        "sheet": sheet.title,
                        "row": row_number,
                        "preserve_chunk": True,
                    },
                )
            )

        if header:
            docs.extend(excel_summary_docs(uploaded_file.name, sheet.title, header, records))
        docs.extend(row_docs)

    workbook.close()
    return docs


def markdown_record_title(block: str, fallback: str) -> str:
    for line in block.splitlines():
        stripped = line.strip()
        if stripped.startswith("## "):
            return stripped.removeprefix("## ").strip()
    return fallback


def extract_rag_style_markdown(text: str, source_name: str) -> list[Document]:
    docs = []
    title = ""
    body = text

    lines = text.splitlines()
    if lines and lines[0].startswith("# "):
        title = lines[0].removeprefix("# ").strip()
        body = "\n".join(lines[1:]).strip()

    blocks = [block.strip() for block in body.split("\n---") if block.strip()]
    if len(blocks) <= 1:
        return []

    sections = [
        markdown_record_title(block, f"Record {index}")
        for index, block in enumerate(blocks, start=1)
    ]
    docs.append(
        Document(
            page_content="\n".join(
                [
                    f"Markdown Summary for {source_name}",
                    f"Title: {title or source_name}",
                    f"Total records: {len(blocks)}",
                    "Records: " + ", ".join(sections),
                ]
            ),
            metadata={
                "source": source_name,
                "section": "Markdown summary",
                "preserve_chunk": True,
            },
        )
    )

    for index, block in enumerate(blocks, start=1):
        section = sections[index - 1]
        content = f"{title}\n\n{block}".strip() if title else block
        docs.append(
            Document(
                page_content=content,
                metadata={
                    "source": source_name,
                    "section": section,
                    "record": index,
                    "preserve_chunk": True,
                },
            )
        )

    return docs


def extract_markdown_table(text: str, source_name: str) -> list[Document]:
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    table_lines = [line for line in lines if line.startswith("|") and line.endswith("|")]

    if len(table_lines) < 3:
        return []

    headers = [cell.strip() for cell in table_lines[0].strip("|").split("|")]
    docs = []

    for row_index, row in enumerate(table_lines[2:], start=1):
        values = [cell.strip() for cell in row.strip("|").split("|")]
        if not any(values):
            continue

        row_text = "\n".join(
            f"- {header}: {value}" for header, value in zip(headers, values)
        )
        name = values[-1] if values else f"Row {row_index}"
        docs.append(
            Document(
                page_content=row_text,
                metadata={
                    "source": source_name,
                    "section": name,
                    "record": row_index,
                    "preserve_chunk": True,
                },
            )
        )

    return docs


def extract_markdown_sections(text: str, source_name: str) -> list[Document]:
    docs = []
    current_heading = None
    current_lines = []

    for line in text.splitlines():
        if line.startswith("## "):
            if current_lines:
                section = current_heading or f"Section {len(docs) + 1}"
                docs.append(
                    Document(
                        page_content="\n".join(current_lines).strip(),
                        metadata={
                            "source": source_name,
                            "section": section,
                            "preserve_chunk": True,
                        },
                    )
                )
            current_heading = line.removeprefix("## ").strip()
            current_lines = [line]
        else:
            current_lines.append(line)

    if current_lines:
        section = current_heading or "Markdown document"
        docs.append(
            Document(
                page_content="\n".join(current_lines).strip(),
                metadata={
                    "source": source_name,
                    "section": section,
                    "preserve_chunk": True,
                },
            )
        )

    return [doc for doc in docs if doc.page_content.strip()]


def extract_markdown(uploaded_file) -> list[Document]:
    text = uploaded_file.getvalue().decode("utf-8", errors="replace").strip()
    if not text:
        return []

    rag_style_docs = extract_rag_style_markdown(text, uploaded_file.name)
    if rag_style_docs:
        return rag_style_docs

    table_docs = extract_markdown_table(text, uploaded_file.name)
    if table_docs:
        return table_docs

    return extract_markdown_sections(text, uploaded_file.name)


def extract_documents(uploaded_files) -> list[Document]:
    documents = []

    for uploaded_file in uploaded_files:
        source_name = uploaded_file.name
        extension = Path(source_name).suffix.lower()

        try:
            if uploaded_file.type == "application/pdf" or extension == ".pdf":
                file_docs = extract_pdf(uploaded_file)
                st.success(f"PDF uploaded: {source_name}")
            elif extension == ".docx":
                file_docs = extract_docx(uploaded_file)
                st.success(f"DOCX uploaded: {source_name}")
            elif extension == ".pptx":
                file_docs = extract_pptx(uploaded_file)
                st.success(f"PPTX uploaded: {source_name}")
            elif extension in {".xlsx", ".xlsm"}:
                file_docs = extract_excel(uploaded_file)
                st.success(f"Excel uploaded: {source_name}")
            elif extension in {".md", ".markdown"}:
                file_docs = extract_markdown(uploaded_file)
                st.success(f"Markdown uploaded: {source_name}")
            else:
                st.warning(f"Unsupported file skipped: {source_name}")
                continue

            documents.extend(file_docs)
        except Exception as exc:
            st.error(f"Failed to process {source_name}: {exc}")

    return [doc for doc in documents if doc.page_content.strip()]


def source_label(doc: Document) -> str:
    source = doc.metadata.get("source", "Unknown")
    if "page" in doc.metadata:
        return f"{source} (page {doc.metadata['page']})"
    if "slide" in doc.metadata:
        return f"{source} (slide {doc.metadata['slide']})"
    if "sheet" in doc.metadata and "row" in doc.metadata:
        return f"{source} (sheet {doc.metadata['sheet']}, row {doc.metadata['row']})"
    if "section" in doc.metadata:
        return f"{source} ({doc.metadata['section']})"
    if "sheet" in doc.metadata:
        return f"{source} (sheet {doc.metadata['sheet']})"
    return source


def format_context(docs: list[Document]) -> str:
    return "\n\n".join(
        f"Source: {source_label(doc)}\nContent:\n{doc.page_content}" for doc in docs
    )


def build_namespace(signature: str) -> str:
    session_id = ensure_session_id()
    return f"rag-{session_id[:12]}-{signature[:16]}"


def clear_previous_namespace(index, namespace: str | None):
    if not namespace:
        return

    try:
        index.delete(delete_all=True, namespace=namespace)
    except Exception:
        pass


uploaded_files = st.file_uploader(
    "Upload Files (PDF, DOCX, PPTX, XLSX, MD)",
    type=["pdf", "docx", "pptx", "xlsx", "xlsm", "md", "markdown"],
    accept_multiple_files=True,
)

if not uploaded_files:
    st.info("Upload one or more PDF, DOCX, PPTX, XLSX, or MD files to start chatting.")
    st.stop()

index = get_pinecone_index(PINECONE_API_KEY, INDEX_NAME)
embedding_model = get_embedding_model()
reranker = get_reranker()
llm = get_llm(GROQ_API_KEY)

signature = uploaded_files_signature(uploaded_files)
namespace = build_namespace(signature)

if st.session_state.get("processed_signature") != signature:
    with st.spinner("Processing uploaded files..."):
        documents = extract_documents(uploaded_files)

        if not documents:
            st.warning("No readable text found in the uploaded files.")
            st.stop()

        preserved_chunks = [
            doc for doc in documents if doc.metadata.get("preserve_chunk")
        ]
        splittable_documents = [
            doc for doc in documents if not doc.metadata.get("preserve_chunk")
        ]

        splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
        split_chunks = splitter.split_documents(splittable_documents)
        chunks = preserved_chunks + split_chunks
        chunks = [chunk for chunk in chunks if chunk.page_content.strip()]

        if not chunks:
            st.warning("No readable chunks could be created from the uploaded files.")
            st.stop()

        previous_namespace = st.session_state.get("namespace")
        if previous_namespace and previous_namespace != namespace:
            clear_previous_namespace(index, previous_namespace)

        vectorstore = PineconeVectorStore(
            index=index,
            embedding=embedding_model,
            namespace=namespace,
        )
        ids = [f"{namespace}-{idx}" for idx in range(len(chunks))]
        vectorstore.add_documents(chunks, ids=ids)

        st.session_state.processed_signature = signature
        st.session_state.namespace = namespace
        st.session_state.chunks = chunks

chunks = st.session_state.get("chunks", [])
if not chunks:
    st.warning("Upload documents before asking a question.")
    st.stop()

st.write(f"Total chunks: {len(chunks)}")

vectorstore = PineconeVectorStore(
    index=index,
    embedding=embedding_model,
    namespace=st.session_state.namespace,
)
retriever = vectorstore.as_retriever(search_kwargs={"k": 5})
bm25 = BM25Okapi([doc.page_content.split() for doc in chunks])


def hybrid_search(query: str, k: int = 5) -> list[Document]:
    tokenized_query = query.split()
    bm25_scores = bm25.get_scores(tokenized_query)
    bm25_results = sorted(
        zip(chunks, bm25_scores),
        key=lambda item: item[1],
        reverse=True,
    )[:k]

    bm25_docs = [doc for doc, _ in bm25_results]
    vector_docs = retriever.invoke(query)[:k]

    seen = set()
    combined = []

    for doc in bm25_docs + vector_docs:
        key = (
            doc.metadata.get("source"),
            doc.metadata.get("page"),
            doc.metadata.get("slide"),
            doc.metadata.get("sheet"),
            doc.metadata.get("row"),
            doc.metadata.get("section"),
            doc.metadata.get("record"),
            doc.page_content,
        )
        if key not in seen:
            combined.append(doc)
            seen.add(key)

    if not combined:
        return []

    pairs = [(query, doc.page_content) for doc in combined]
    scores = reranker.predict(pairs)
    reranked = sorted(zip(combined, scores), key=lambda item: item[1], reverse=True)

    return [doc for doc, _ in reranked[:3]]


def normalize_query_text(text: str) -> str:
    return re.sub(r"[^a-z0-9]+", " ", text.lower()).strip()


def line_value(text: str, label: str) -> str | None:
    pattern = rf"^{re.escape(label)}:\s*(.+)$"
    match = re.search(pattern, text, flags=re.MULTILINE | re.IGNORECASE)
    return match.group(1).strip() if match else None


def split_csv_value(value: str | None) -> list[str]:
    if not value:
        return []
    return [item.strip() for item in value.split(",") if item.strip()]


def parse_row_document(doc: Document) -> dict[str, str] | None:
    if "sheet" not in doc.metadata or "row" not in doc.metadata:
        return None

    record = {
        "Source": str(doc.metadata.get("source", "")),
        "Sheet": str(doc.metadata.get("sheet", "")),
        "Row": str(doc.metadata.get("row", "")),
    }

    for line in doc.page_content.splitlines():
        stripped = line.strip()
        if not stripped.startswith("- ") or ": " not in stripped:
            continue

        key, value = stripped[2:].split(": ", 1)
        record[key.strip()] = value.strip()

    return record if len(record) > 3 else None


def structured_rows_from_chunks(available_chunks: list[Document]) -> list[dict[str, str]]:
    rows = []

    for doc in available_chunks:
        row = parse_row_document(doc)
        if row:
            rows.append(row)

    return rows


def is_vm_table(rows: list[dict[str, str]]) -> bool:
    if not rows:
        return False

    keys = {key.lower() for row in rows[:20] for key in row}
    return bool({"vm name", "vm ip", "service name", "port", "client"} & keys)


def looks_like_structured_vm_question(query: str, rows: list[dict[str, str]]) -> bool:
    if not is_vm_table(rows):
        return False

    normalized_query = normalize_query_text(query)
    field_terms = {
        "client",
        "nbp",
        "project",
        "environment",
        "env",
        "vm",
        "vms",
        "machine",
        "machines",
        "ip",
        "service",
        "services",
        "port",
        "ports",
        "raast",
        "mpay",
        "digital",
        "jazz",
        "beoe",
        "bulk",
    }

    if any(term in normalized_query.split() for term in field_terms):
        return True

    for row in rows[:200]:
        for key in ["VM Name", "VM IP", "Service Name", "Project", "Environment"]:
            value = normalize_query_text(row.get(key, ""))
            if value and value in normalized_query:
                return True

    return False


def compact_structured_table(rows: list[dict[str, str]], max_rows: int = 350) -> str:
    preferred_columns = [
        "Source",
        "Sheet",
        "Row",
        "S.No",
        "Client",
        "Project",
        "Environment",
        "VM IP",
        "VM Name",
        "Service Name",
        "Port",
        "Purpose",
    ]
    all_columns = ordered_unique([key for row in rows for key in row.keys()])
    columns = [column for column in preferred_columns if column in all_columns]
    columns.extend(column for column in all_columns if column not in columns)

    lines = [" | ".join(columns)]
    lines.append(" | ".join("---" for _ in columns))

    for row in rows[:max_rows]:
        lines.append(" | ".join(row.get(column, "") for column in columns))

    if len(rows) > max_rows:
        lines.append(f"... {len(rows) - max_rows} more rows omitted ...")

    return "\n".join(lines)


def structured_table_summary(rows: list[dict[str, str]]) -> str:
    if not rows:
        return ""

    clients = ordered_unique([row.get("Client", "") for row in rows])
    unique_vms = ordered_unique([row.get("VM Name", "") for row in rows])
    unique_services = ordered_unique([row.get("Service Name", "") for row in rows])
    projects = count_values(rows, "Project") if "Project" in rows[0] else {}
    environments = count_values(rows, "Environment") if "Environment" in rows[0] else {}

    lines = [
        f"Total rows: {len(rows)}",
        f"Clients: {', '.join(clients)}",
        f"Total unique VM Names: {len(unique_vms)}",
        "Unique VM Names: " + ", ".join(unique_vms),
    ]

    if unique_services:
        lines.extend(
            [
                f"Total unique Service Names: {len(unique_services)}",
                "Unique Service Names: " + ", ".join(unique_services),
            ]
        )

    if projects:
        lines.append(
            "Project counts: "
            + ", ".join(f"{key}: {value}" for key, value in sorted(projects.items()))
        )
    if environments:
        lines.append(
            "Environment counts: "
            + ", ".join(f"{key}: {value}" for key, value in sorted(environments.items()))
        )

    return "\n".join(lines)


def structured_answer(query: str, available_chunks: list[Document]):
    normalized_query = normalize_query_text(query)
    is_count_query = any(
        word in normalized_query
        for word in ["kitni", "kitna", "count", "total", "how many"]
    )
    is_list_query = any(
        word in normalized_query
        for word in ["kya kya", "list", "names", "naam", "which"]
    )
    wants_vm = any(
        word in normalized_query
        for word in ["vm", "vms", "machine", "machines"]
    )
    wants_service = any(
        word in normalized_query
        for word in ["service", "services"]
    )

    if not (is_count_query or is_list_query) or not (wants_vm or wants_service):
        return None, []

    client_summaries = [
        doc
        for doc in available_chunks
        if str(doc.metadata.get("section", "")).lower().startswith("client summary:")
    ]

    matching_summaries = [
        doc
        for doc in client_summaries
        if normalize_query_text(str(doc.metadata.get("client", ""))) in normalized_query
    ]

    if not matching_summaries and len(client_summaries) == 1:
        matching_summaries = client_summaries

    if not matching_summaries:
        return None, []

    doc = matching_summaries[0]
    client = doc.metadata.get("client", "selected client")

    if wants_vm:
        count = line_value(doc.page_content, "Total unique VMs for client")
        names = split_csv_value(line_value(doc.page_content, "Unique VM Names"))
        if count:
            answer = f"{client} main total unique VMs/machines {count} hain."
            if is_list_query and names:
                answer += "\n\nVM names:\n" + "\n".join(f"- {name}" for name in names)
            return answer, [doc]

    if wants_service:
        count = line_value(doc.page_content, "Total unique services for client")
        names = split_csv_value(line_value(doc.page_content, "Unique Services"))
        if count:
            answer = f"{client} main total unique services {count} hain."
            if is_list_query and names:
                answer += "\n\nServices:\n" + "\n".join(f"- {name}" for name in names)
            return answer, [doc]

    return None, []


def structured_vm_answer(query: str, available_chunks: list[Document], llm):
    rows = structured_rows_from_chunks(available_chunks)
    if not looks_like_structured_vm_question(query, rows):
        return None, []

    context = "\n\n".join(
        [
            "Structured VM/service data summary:",
            structured_table_summary(rows),
            "Structured VM/service rows:",
            compact_structured_table(rows),
        ]
    )

    prompt = f"""
You answer questions about the uploaded VM/service inventory.

Use ONLY the structured data below. Do not guess.
For questions about "VMs", "machines", or "servers", count DISTINCT VM Name values unless the user explicitly asks for rows, services, ports, or entries.
For questions about services, use Service Name.
If the answer is not present in the structured data, say "I don't know from the uploaded VM list."
Give concise answers, and include relevant VM names, IPs, services, ports, project, or environment when useful.

{context}

Question:
{query}
"""

    response = llm.invoke(prompt)
    sources = [
        doc
        for doc in available_chunks
        if str(doc.metadata.get("section", "")).lower().startswith("excel summary")
        or str(doc.metadata.get("section", "")).lower().startswith("client summary:")
    ]
    return response.content, sources[:3]


query = st.text_input("Ask a question:").strip()

if query:
    with st.spinner("Generating answer..."):
        direct_answer, direct_sources = structured_answer(query, chunks)
        if not direct_answer:
            direct_answer, direct_sources = structured_vm_answer(query, chunks, llm)

        if direct_answer:
            st.write("### Answer")
            st.write(direct_answer)

            st.write("### Sources")
            for source in sorted({source_label(doc) for doc in direct_sources}):
                st.write(f"- {source}")

            docs = []
        else:
            docs = hybrid_search(query)

        if not direct_answer and not docs:
            st.warning("No relevant info found.")
        elif not direct_answer:
            context = format_context(docs)
            sources = sorted({source_label(doc) for doc in docs})

            prompt = f"""
You are a helpful assistant.

Use ONLY the provided context.
If the answer is not found in the context, say "I don't know".

Context:
{context}

Question:
{query}
"""

            response = llm.invoke(prompt)

            st.write("### Answer")
            st.write(response.content)

            st.write("### Sources")
            for source in sources:
                st.write(f"- {source}")


